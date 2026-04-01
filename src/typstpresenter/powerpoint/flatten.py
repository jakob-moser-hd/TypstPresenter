from collections.abc import Generator

from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.shapes.shapetree import _BaseShapes


def flatten(shapes: _BaseShapes) -> Generator[BaseShape]:
    """
    Flatten a shape tree into a single iterable over all shapes in it.
    """

    for shape_or_group in shapes:
        if isinstance(shape_or_group, GroupShape):
            yield from flatten(shape_or_group.shapes)
        else:
            yield shape_or_group
