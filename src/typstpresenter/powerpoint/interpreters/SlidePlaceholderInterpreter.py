import logging
from itertools import groupby

from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from pptx.shapes import Subshape
from pptx.shapes.base import BaseShape
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.text.text import TextFrame, _Paragraph, _Run

from typstpresenter.model.Element import Element
from typstpresenter.model.List import List
from typstpresenter.model.PresentationTitle import PresentationTitle
from typstpresenter.model.Title import Title
from typstpresenter.model.text.Atom import Atom
from typstpresenter.model.text.Link import Link
from typstpresenter.model.text.Subscript import Subscript
from typstpresenter.model.text.Superscript import Superscript
from typstpresenter.model.text.Text import Text
from typstpresenter.powerpoint.Ignore import Ignore


type Level = int


logger = logging.getLogger(__name__)


class SlidePlaceholderInterpreter:
    """
    Interpret a shape of type SlidePlaceholder.

    These are somewhat nicer to interpret, because some of them have a type assigned which is actually semantic.

    See: https://python-pptx.readthedocs.io/en/latest/user/understanding-shapes.html
    """
    def can_interpret(self, shape: BaseShape | Subshape) -> bool:
        return isinstance(shape, SlidePlaceholder)

    def __call__(self, shape: BaseShape | Subshape, context: dict | None = None) -> Element | Ignore | None:
        if not isinstance(shape, SlidePlaceholder):
            return None

        match shape.placeholder_format.type:
            case PP_PLACEHOLDER_TYPE.TITLE:
                return Title(text=_interpret_text_frame(shape.text_frame, default_to_list=False))
            case PP_PLACEHOLDER_TYPE.CENTER_TITLE:
                return PresentationTitle(text=_interpret_text_frame(shape.text_frame, default_to_list=False))
            case PP_PLACEHOLDER_TYPE.SLIDE_NUMBER | PP_PLACEHOLDER_TYPE.HEADER | PP_PLACEHOLDER_TYPE.FOOTER | PP_PLACEHOLDER_TYPE.DATE:
                return Ignore()
            case PP_PLACEHOLDER_TYPE.SUBTITLE:
                return _interpret_text_frame(shape.text_frame, default_to_list=False)
            case PP_PLACEHOLDER_TYPE.OBJECT:
                # Just pretend that object means a bunch of text, and nothing else.
                return _interpret_text_frame(shape.text_frame, default_to_list=True)
            case _:
                return None


def _interpret_text_frame(text_frame: TextFrame, default_to_list: bool = True) -> Text | List:
    if len(text_frame.paragraphs) == 1:
        return _interpret_paragraph(text_frame.paragraphs[0])

    if not default_to_list:
        atoms = []
        for p in text_frame.paragraphs:
            atoms.extend(_interpret_paragraph(p).value)
            atoms.append("\n")
        if atoms:
            atoms.pop()
        return Text(tuple(atoms))

    # Just pretend any multi-paragraph text is a list
    paragraphs_by_level = groupby(text_frame.paragraphs, key=lambda item: item.level)
    list_stack: list[tuple[Level, List]] = []

    for level, paragraphs in paragraphs_by_level:
        list_ = List(tuple(_interpret_paragraph(p) for p in paragraphs))

        if len(list_stack) > 0:
            previous_level, previous_list = list_stack.pop()
        else:
            # Stack was previously empty, this only happens in the first iteration.
            previous_level, previous_list = -1, None

        # Go up the stack until we find a parent (i.e. one whose level is smaller than the current one).
        # If the current list is a root level, we'd empty the stack completely
        while len(list_stack) > 0 and level <= previous_level:
            previous_level, previous_list = list_stack.pop()

        if level > previous_level:
            if previous_list is not None:
                # This list has a parent, append it to the parent and put the barent back on the stack.
                list_stack.append((previous_level, previous_list.append(list_)))

            # Put the list itself on the stack because it is the level we are currently working at.
            list_stack.append((level, list_))
        elif level == previous_level:
            # If this list is on the same level as the previous one, concatenate both lists and keep the
            # earlier one.
            list_stack.append((previous_level, previous_list.append(*list_)))

    if list_stack:
        root_level, root_list = list_stack[0]

        return root_list
    else:
        logger.warning("Found ill-formed list, returning nothing")
        return Text(("",))


def _interpret_paragraph(paragraph: _Paragraph) -> Text:
    return Text(tuple(_interpret_run(run) for run in paragraph.runs))


# See https://stackoverflow.com/questions/61329224/how-do-i-add-superscript-subscript-text-to-powerpoint-using-python-pptx
_SUBSCRIPT = "-25000"
_SUPERSCRIPT = "30000"


def _interpret_run(run: _Run) -> Atom:
    if run.hyperlink.address is not None:
        return Link(text=Text(run.text), target=run.hyperlink.address)

    # The text baseline indicates how high in a line the run is positioned.
    baseline_position = run.font._element.get("baseline")

    if baseline_position == _SUBSCRIPT:
        return Subscript(text=Text(run.text))
    elif baseline_position == _SUPERSCRIPT:
        return Superscript(text=Text(run.text))

    return run.text
