from pptx.shapes import Subshape
from pptx.shapes.base import BaseShape
from pptx.shapes.placeholder import SlidePlaceholder

from typstpresenter.model.Element import Element
from typstpresenter.powerpoint.Ignore import Ignore
from typstpresenter.powerpoint.interpreters.SlidePlaceholderInterpreter import _interpret_text_frame


class TextBoxInterpreter:
    """
    Interpret generic shapes that contain a text frame (e.g. manually added text boxes).
    """
    def can_interpret(self, shape: BaseShape | Subshape) -> bool:
        return hasattr(shape, "text_frame") and not isinstance(shape, SlidePlaceholder)

    def __call__(self, shape: BaseShape | Subshape, context: dict | None = None) -> Element | Ignore | None:
        if hasattr(shape, "text_frame") and getattr(shape, "has_text_frame", False):
            return _interpret_text_frame(shape.text_frame, default_to_list=False)
        return None
