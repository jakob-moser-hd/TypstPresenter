from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes import Subshape
from pptx.shapes.base import BaseShape

from typstpresenter.model.Element import Element
from typstpresenter.model.Image import Image
from typstpresenter.powerpoint.Ignore import Ignore


class PictureInterpreter:
    def can_interpret(self, shape: BaseShape | Subshape) -> bool:
        return hasattr(shape, "shape_type") and shape.shape_type == MSO_SHAPE_TYPE.PICTURE

    def __call__(self, shape: BaseShape | Subshape, context: dict | None = None) -> Element | Ignore | None:
        context = context or {}
        slide_idx = context.get("slide_index", 0)
        elem_idx = context.get("element_index", 0)

        ext = shape.image.ext
        name = f"slide_{slide_idx + 1}_pos_{elem_idx}.{ext}"
        blob = shape.image.blob
        width = getattr(shape, 'width', None)
        height = getattr(shape, 'height', None)
        width_pt = getattr(width, 'pt', None) if width else None
        height_pt = getattr(height, 'pt', None) if height else None
        return Image(name=name, ext=ext, blob=blob, width_pt=width_pt, height_pt=height_pt)
