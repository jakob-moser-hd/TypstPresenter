from collections.abc import Sequence
from dataclasses import dataclass
from pathlib import Path
from typing import Self, overload

import pptx

from typstpresenter.model.PresentationTitle import PresentationTitle
from typstpresenter.model.Slide import Slide
from typstpresenter.typst.express import express


@dataclass
class Presentation(Sequence[Slide]):
    title: PresentationTitle
    slides: Sequence[Slide]

    # The path to the source file (if sourced from a file)
    source_path: Path | None = None

    @classmethod
    def from_file(cls, path: Path) -> Self:
        """
        Load a presentation from the file at the given path.

        Currently assumes that the path points to a *.pptx file, no other file types can be handled.
        Will fail if other files are presented, possibly in curious ways.
        """
        if not path.is_file():
            raise FileNotFoundError(f"'{path}' is not a file")

        prs = pptx.Presentation(str(path))
        slides = tuple(Slide.from_pptx_slide(pptx_slide, slide_index=i) for i, pptx_slide in enumerate(prs.slides))

        # The title slide is metadata in my world, so we remove it from the set of slides later, and extract the
        # metadata here in a dedicated way.
        title_slide = slides[0]

        return cls(
            title=title_slide.get_singular_element_or_none(PresentationTitle),
            slides=slides[1:],
            source_path=path,
        )

    def to_typst_str(self, media_dir: str = "media") -> str:
        """
        Convert the presentation to a string in Typst format and return it.
        """
        return express(self, media_dir=media_dir)

    def to_file(self, path: Path) -> None:
        """
        Save a presentation to the given path.

        Will output the file in *.typ (Typst) format, no matter the extension.
        Extracts any media required (e.g., images) to a neighboring folder.
        """
        media_dir_name = f"{path.stem}_media"
        media_dir_path = path.parent / media_dir_name

        from typstpresenter.model.Image import Image
        images = []
        for slide in self.slides:
            for placed in slide.elements:
                if isinstance(placed.element, Image):
                    images.append(placed.element)

        if images:
            media_dir_path.mkdir(exist_ok=True, parents=True)
            for img in images:
                (media_dir_path / img.name).write_bytes(img.blob)

        with path.open("w", encoding="utf-8") as f:
            f.write(self.to_typst_str(media_dir=media_dir_name))

    @overload
    def __getitem__(self, index: int) -> Slide: ...

    @overload
    def __getitem__(self, index: slice) -> Sequence[Slide]: ...

    def __getitem__(self, index):
        return self.slides[index]

    def __len__(self) -> int:
        return len(self.slides)
