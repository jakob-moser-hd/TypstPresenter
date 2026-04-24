import pptx
from pptx.util import Inches
from pathlib import Path
from dataclasses import dataclass


@dataclass
class PresentationConfig:
    """Konfiguration für die Erstellung der Test-Präsentationen."""
    slide_width_in: float = 10.0
    slide_height_in: float = 7.5
    margin_x_in: float = 1.0
    margin_y_in: float = 1.5   # Etwas mehr Platz oben/unten
    gutter_in: float = 0.5


def add_title_slide(prs: pptx.Presentation, title_text: str, subtitle_text: str):
    """Fügt der Präsentation eine standardisierte Titelfolie hinzu."""
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text


def add_slide_with_title(prs: pptx.Presentation, title_text: str):
    """Fügt eine Folie hinzu, die nur einen Titel besitzt (Layout 5), und gibt diese zurück."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title_text
    return slide


def generate_placeholder_text(description: str, item_count: int = 1) -> str:
    """Generiert einen längeren Text mit einer führenden Beschreibung für Test-Placeholder.
    Die Länge des Textes skaliert invers zur Anzahl der Elemente auf der Folie."""
    lorem_ipsum = (
        "Lorem ipsum dolor sit amet, consectetur adipiscing elit. "
        "Aenean commodo ligula eget dolor. Aenean massa. Cum sociis "
        "natoque penatibus et magnis dis parturient montes, nascetur "
        "ridiculus mus. Donec quam felis, ultricies nec, pellentesque eu."
    )
    sentences = lorem_ipsum.split(". ")
    num_sentences = max(1, 5 - item_count)
    shortened_text = ". ".join(sentences[:num_sentences])
    if not shortened_text.endswith("."):
        shortened_text += "."
        
    return f"[{description}]\n\n{shortened_text}"


def create_simple_pptx(path: Path, config: PresentationConfig):
    prs = pptx.Presentation()
    
    # Title Slide
    add_title_slide(
        prs,
        title_text="Hello World Presentation",
        subtitle_text="Eine einfache Präsentation zu Testzwecken. Sie enthält eine Titelfolie, "
                      "um grundlegende Konvertierungen und Metadaten-Extraktionen prüfen zu "
                      "können, und illustriert ein einfaches Bullet-Point Layout."
    )
    
    # Bullet Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Slide 1 Title"
    
    tf = body_shape.text_frame
    tf.text = "This is a bullet point"
    
    p = tf.add_paragraph()
    p.text = "This is a second bullet point"
    
    prs.save(path)


def create_two_content_pptx(path: Path, config: PresentationConfig):
    prs = pptx.Presentation()
    
    # Title Slide
    add_title_slide(
        prs,
        title_text="Two Content Presentation",
        subtitle_text="Diese Testpräsentation prüft das 'Two Content' Layout. "
                      "Dabei ist wichtig, dass beide Spalten als Layout Element korrekt in das Grid "
                      "im Typst-Code extrahiert und gesetzt werden."
    )
    
    # "Two Content" slide is usually layout index 3 in the default template
    slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    title_shape.text = "Two Content Slide"
    
    body_shape_1 = shapes.placeholders[1]
    body_shape_2 = shapes.placeholders[2]
    
    body_shape_1.text_frame.text = generate_placeholder_text("Two content layout, left block.", item_count=2)
    body_shape_2.text_frame.text = generate_placeholder_text("Two content layout, right block.", item_count=2)
    
    prs.save(path)


def create_multi_content_pptx(path: Path, config: PresentationConfig):
    prs = pptx.Presentation()
    
    # Title Slide
    add_title_slide(
        prs,
        title_text="Multi Content Presentation",
        subtitle_text="Eine komplexere Präsentation, um mehrzeilige  "
                      "und mehrspaltige Layouts zu prüfen."
    )
    
    # Berechne den nutzbaren Platz basierend auf der Konfiguration
    usable_width = config.slide_width_in - 2 * config.margin_x_in
    usable_height = config.slide_height_in - 2 * config.margin_y_in
    
    def add_textbox(slide, x_in: float, y_in: float, w_in: float, h_in: float, text: str):
        """Hilfsfunktion zum Einfügen von flexiblen Textboxen."""
        tx = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
        tx.text_frame.word_wrap = True  # Zeilenumbruch aktivieren
        tx.text_frame.text = text
        return tx

    # --- 1. Slide: two rows (1 placeholder each) ---
    slide1 = add_slide_with_title(prs, "Two Rows Layout")
    row_height_2 = (usable_height - config.gutter_in) / 2
    
    add_textbox(slide1, config.margin_x_in, config.margin_y_in, usable_width, row_height_2,
                generate_placeholder_text("Two rows setup, top row (1/2).", item_count=2))
    add_textbox(slide1, config.margin_x_in, config.margin_y_in + row_height_2 + config.gutter_in, usable_width, row_height_2,
                generate_placeholder_text("Two rows setup, bottom row (2/2).", item_count=2))

    # --- 2. Slide: three rows ---
    slide2 = add_slide_with_title(prs, "Three Rows Layout")
    row_height_3 = (usable_height - 2 * config.gutter_in) / 3
    
    for i in range(3):
        y_pos = config.margin_y_in + i * (row_height_3 + config.gutter_in)
        add_textbox(slide2, config.margin_x_in, y_pos, usable_width, row_height_3,
                    generate_placeholder_text(f"Three rows setup, row {i+1} of 3.", item_count=3))

    # --- 3. Slide: a 2x2 grid of placeholders, all 4 present ---
    slide3 = add_slide_with_title(prs, "Full 2x2 Grid")
    col_width_2 = (usable_width - config.gutter_in) / 2
    
    positions_2x2 = [
        (config.margin_x_in, config.margin_y_in),                                                          # TL
        (config.margin_x_in + col_width_2 + config.gutter_in, config.margin_y_in),                         # TR
        (config.margin_x_in, config.margin_y_in + row_height_2 + config.gutter_in),                        # BL
        (config.margin_x_in + col_width_2 + config.gutter_in, config.margin_y_in + row_height_2 + config.gutter_in)  # BR
    ]
    labels_2x2 = ["top-left", "top-right", "bottom-left", "bottom-right"]

    for i, ((x, y), label) in enumerate(zip(positions_2x2, labels_2x2)):
        add_textbox(slide3, x, y, col_width_2, row_height_2,
                    generate_placeholder_text(f"Full 2x2 grid, placeholder {i+1} of 4 ({label}).", item_count=4))

    # --- 4. Slides: sequence of 1, 2, 3, 4 placeholders present in a 2x2 grid ---
    for count in range(1, 5):
        s = add_slide_with_title(prs, f"Sequential Grid ({count}/4)")
        for i in range(count):
            x, y = positions_2x2[i]
            label = labels_2x2[i]
            add_textbox(s, x, y, col_width_2, row_height_2,
                        generate_placeholder_text(f"Sequential 2x2 grid (Slide {count}/4) - Placeholder {i+1} ({label}).\nHat {count} von 4 Platzhaltern.", item_count=count))
            
    prs.save(path)


def get_dummy_image():
    import base64
    import io
    # Basic 1x1 red pixel PNG
    img_b64 = b"iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
    return io.BytesIO(base64.b64decode(img_b64))


def create_media_pptx(path: Path, config: PresentationConfig):
    prs = pptx.Presentation()
    add_title_slide(prs, "Media Layouts Test", "Slides with Images alongside Text")
    
    # Slide 1: Image Left, Text Right
    slide = add_slide_with_title(prs, "Image and Text")
    
    content_width = (config.slide_width_in - 2 * config.margin_x_in - config.gutter_in) / 2
    content_height = config.slide_height_in - config.margin_y_in - 1.0
    
    slide.shapes.add_picture(
        get_dummy_image(), 
        Inches(config.margin_x_in), 
        Inches(config.margin_y_in), 
        width=Inches(content_width), 
        height=Inches(content_height)
    )
    
    tx = slide.shapes.add_textbox(
        Inches(config.margin_x_in + content_width + config.gutter_in), 
        Inches(config.margin_y_in), 
        Inches(content_width), 
        Inches(content_height)
    )
    tx.text_frame.word_wrap = True
    tx.text_frame.text = generate_placeholder_text("Text next to an image.", item_count=2)
    
    prs.save(path)


if __name__ == "__main__":
    out_dir = Path("tests/data")
    out_dir.mkdir(exist_ok=True, parents=True)
    
    config = PresentationConfig()
    
    create_simple_pptx(out_dir / "simple.pptx", config)
    create_two_content_pptx(out_dir / "two_content.pptx", config)
    create_multi_content_pptx(out_dir / "multi_content.pptx", config)
    create_media_pptx(out_dir / "media.pptx", config)
    
    print("Testdaten wurden erfolgreich generiert!")
